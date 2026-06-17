from pathlib import Path

from fastapi import HTTPException, status
from langchain_community.document_loaders import PyPDFLoader
from langchain_core.documents import Document

OCR_ERROR = "Для распознавания текста с изображения нужно установить OCR-модуль."


def extract_documents(file_path: str, extension: str) -> list[Document]:
    path = Path(file_path)
    if not path.exists():
        raise HTTPException(
            status_code=status.HTTP_404_NOT_FOUND,
            detail="Загруженный файл не найден на диске.",
        )

    if extension == ".pdf":
        documents = PyPDFLoader(str(path)).load()
    elif extension in {".txt", ".md"}:
        documents = [_load_txt(path)]
    elif extension == ".docx":
        documents = [_load_docx(path)]
    elif extension == ".pptx":
        documents = _load_pptx(path)
    elif extension == ".xlsx":
        documents = _load_xlsx(path)
    elif extension in {".png", ".jpg", ".jpeg"}:
        documents = [_load_image(path)]
    else:
        raise HTTPException(
            status_code=status.HTTP_400_BAD_REQUEST,
            detail="Неподдерживаемый формат файла. Поддерживаемые форматы: PDF, DOCX, TXT, MD, PPTX, XLSX, PNG, JPG.",
        )

    documents = [doc for doc in documents if clean_text(doc.page_content)]
    if not documents:
        raise HTTPException(
            status_code=status.HTTP_400_BAD_REQUEST,
            detail="Не удалось извлечь текст из файла.",
        )
    return documents


def clean_text(text: str) -> str:
    lines = [" ".join(line.split()) for line in text.splitlines()]
    return "\n".join(line for line in lines if line).strip()


def _load_txt(path: Path) -> Document:
    for encoding in ("utf-8", "cp1251"):
        try:
            text = path.read_text(encoding=encoding)
            return Document(
                page_content=text,
                metadata={"source": str(path), "page": None, "encoding": encoding},
            )
        except UnicodeDecodeError:
            continue
    raise HTTPException(
        status_code=status.HTTP_400_BAD_REQUEST,
        detail="Не удалось прочитать текстовый файл в кодировке UTF-8 или CP1251.",
    )


def _load_docx(path: Path) -> Document:
    try:
        from docx import Document as DocxDocument
    except ImportError as exc:
        raise HTTPException(
            status_code=status.HTTP_500_INTERNAL_SERVER_ERROR,
            detail="Для чтения DOCX нужно установить python-docx.",
        ) from exc

    doc = DocxDocument(str(path))
    paragraphs = [paragraph.text for paragraph in doc.paragraphs if paragraph.text.strip()]
    table_rows = []
    for table in doc.tables:
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
            if cells:
                table_rows.append(" | ".join(cells))
    return Document(
        page_content="\n".join(paragraphs + table_rows),
        metadata={"source": str(path), "page": None},
    )


def _load_pptx(path: Path) -> list[Document]:
    try:
        from pptx import Presentation
    except ImportError as exc:
        raise HTTPException(
            status_code=status.HTTP_500_INTERNAL_SERVER_ERROR,
            detail="Для чтения PPTX нужно установить python-pptx.",
        ) from exc

    presentation = Presentation(str(path))
    documents = []
    for slide_index, slide in enumerate(presentation.slides, start=1):
        texts = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                texts.append(shape.text.strip())
        if texts:
            documents.append(
                Document(
                    page_content="\n".join(texts),
                    metadata={"source": str(path), "page": slide_index},
                )
            )
    return documents


def _load_xlsx(path: Path) -> list[Document]:
    try:
        from openpyxl import load_workbook
    except ImportError as exc:
        raise HTTPException(
            status_code=status.HTTP_500_INTERNAL_SERVER_ERROR,
            detail="Для чтения XLSX нужно установить openpyxl.",
        ) from exc

    workbook = load_workbook(str(path), read_only=True, data_only=True)
    documents = []
    for sheet_index, sheet in enumerate(workbook.worksheets, start=1):
        rows = []
        for row in sheet.iter_rows(values_only=True):
            values = [str(value).strip() for value in row if value is not None and str(value).strip()]
            if values:
                rows.append(" | ".join(values))
        if rows:
            documents.append(
                Document(
                    page_content=f"Лист: {sheet.title}\n" + "\n".join(rows),
                    metadata={"source": str(path), "page": sheet_index, "sheet": sheet.title},
                )
            )
    workbook.close()
    return documents


def _load_image(path: Path) -> Document:
    try:
        from PIL import Image
        import pytesseract
    except ImportError as exc:
        raise HTTPException(status_code=status.HTTP_400_BAD_REQUEST, detail=OCR_ERROR) from exc

    try:
        text = pytesseract.image_to_string(Image.open(path), lang="rus+eng")
    except pytesseract.TesseractNotFoundError as exc:
        raise HTTPException(status_code=status.HTTP_400_BAD_REQUEST, detail=OCR_ERROR) from exc
    except Exception as exc:
        raise HTTPException(
            status_code=status.HTTP_400_BAD_REQUEST,
            detail="Не удалось распознать текст на изображении.",
        ) from exc
    return Document(page_content=text, metadata={"source": str(path), "page": None})
